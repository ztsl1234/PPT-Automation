import logging
import os
from os import path

#from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches
#import argparse
import pandas as pd
import numpy as np
from datetime import date
import matplotlib.pyplot as plt
#import seaborn as sns
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

'''
""" Ref for slide types/Layout: 
0 ->  title and subtitle
1 ->  title and content
2 ->  section header
3 ->  two content
4 ->  Comparison
5 ->  Title only 
6 ->  Blank
7 ->  Content with caption
8 ->  Pic with caption
"""
template 
phf_idx,shape name,phf_type
0, Title 1, CENTER_TITLE (3)
1, Subtitle 2, SUBTITLE (4)
0, Title 1, TITLE (1)
13, Content Placeholder 2, OBJECT (7)
18, Table Placeholder 3, TABLE (12)
19, Table Placeholder 4, TABLE (12)
0, Title 2, TITLE (1)
15, Content Placeholder 1, OBJECT (7)
0, Title 1, TITLE (1)
18, Table Placeholder 2, TABLE (12)
19, Table Placeholder 3, TABLE (12)


'''
class Q2c:
    #beige=253,234,218
    #light blue=217,241,255
    #white=252,252,252
    RGB_biege=[0xfd,0xea,0xda]
    RGB_light_blue=[0xd9,0xf1,0xff]
    RGB_white=[0xfc,0xfc,0xfc]    

    def create_chart(self,df, filename):

        my_data = df["MeasureValue"]
        my_labels = df["Type"]
        plt.pie(my_data,labels=my_labels,autopct='%1.1f%%')
        plt.legend(title="level",loc="lower right")
        plt.title("Excitement levels")
        plt.savefig(filename)

        #plt.style.use('ggplot')
        #pie = df.plot(kind="pie", figsize=(6,6), legend = False, use_index=False, subplots=True, colormap="Pastel1")

        #fig = pie.get_figure()
        #fig.savefig(filename)

    def create_ppt(self,input, output, report_data, chart):

        #create ppt using template file
        prs = Presentation(input)

        logging.info("---> generating slide1")
        #slide1
        title_txt    = "My First PPT Automation"
        subtitle_txt = "I DID IT!"
        self.create_slide_layout0(prs,title_txt,subtitle_txt)

        #data=report_data[report_data["PeriodKey"]==1 & report_data["Category"]=="Technical Ability Test" & report_data["Subcategory"]=="PPT Automation" & report_data["Country"]=="Singapore" & report_data["Segment"]=="All" &  report_data["Section"]=="Base"]
        data=report_data[report_data["PeriodKey"]==1]
        data=data[data["Category"]=="Technical Ability Test"]
        data=data[data["Subcategory"]=="PPT Automation"]

        data_sg=data[data["Country"]=="Singapore"]
      
        data_all=data_sg[data_sg["Segment"]=="All"]

        logging.info("---> generating slide2")
        #slide2 
        data_base=data_all[data_all["Section"]=="Base"]
        data_source=data_all[data_all["Section"]=="Source"]
        data_base.sort_values(by='MeasureValue', ascending=True) 

        self.create_slide_layout1(prs,data_base,data_source)

        logging.info("---> generating slide3")
        #slide3
        data_b4_gender=data_sg[data_sg["Segment"]=="Before"]
        data_b4_gender=data_b4_gender[data_b4_gender["Section"]=="Gender Split"]

        data_after_gender=data_sg[data_sg["Segment"]=="After"]
        data_after_gender=data_after_gender[data_after_gender["Section"]=="Gender Split"]

        data_b4_industry=data_sg[data_sg["Segment"]=="Before"]
        data_b4_industry=data_b4_industry[data_b4_industry["Section"]=="Industry"]

        data_after_industry=data_sg[data_sg["Segment"]=="After"]
        data_after_industry=data_after_industry[data_after_industry["Section"]=="Industry"]

        data_b4_gender.sort_values(by='MeasureValue', ascending=True) 
        data_after_gender.sort_values(by='MeasureValue', ascending=True) 
        data_b4_industry.sort_values(by='MeasureValue', ascending=True) 
        data_after_industry.sort_values(by='MeasureValue', ascending=True) 
        self.create_slide_layout3(prs,data_b4_gender,data_after_gender,data_b4_industry,data_after_industry)

        logging.info("---> generating slide4")
        data_excite=data_all[data_all["Section"]=="Excitement levels"]
        self.create_slide_chart(prs,data_excite,chart)

        logging.info("---> saved to file")
        prs.save(output)

    def create_slide_layout0(self,prs,title_txt,subtitle_txt):
        ######slide : title slide - layout 0
        slide_layout_title = prs.slide_layouts[0]

        slide = prs.slides.add_slide(slide_layout_title)

        #get ph in template
        title = slide.shapes.title
        subtitle_ph = slide.placeholders[1]

        #fill ph with content
        title.text = title_txt
        subtitle_ph.text = subtitle_txt

    def create_slide_layout1(self,prs,data_df1,data_df2):
        logging.debug(f" ---> data_df1={data_df1}")    
       
        ######slide : title and content slide - layout 1
        slide_layout_slide2 = prs.slide_layouts[1]

        slide   = prs.slides.add_slide(slide_layout_slide2)

        #get ph in template
        title   = slide.shapes.title
        subtitle_ph  = slide.placeholders[13]
        table1_ph    = slide.placeholders[18]
        table2_ph    = slide.placeholders[19]    

        #fill ph with content
        title.text = "Creating my first slide with tables"
        subtitle_ph.text = "With colours!"

        num_rows=len(data_df1)+1
        num_cols=4
        graphic_frame = table1_ph.insert_table(rows=num_rows, cols=num_cols)
        table1 = graphic_frame.table

        self.fill_table_colour(table1,0,num_rows, num_cols,self.RGB_white)
        #first row
        self.fill_table_colour(table1,0,1, num_cols,self.RGB_biege)  

        for row in range(num_rows):                
            cell_start  = table1.cell(row, 0)
            cell_end    = table1.cell(row, 2)
            cell_start.merge(cell_end)

        num_rows=len(data_df2)+1
        graphic_frame = table2_ph.insert_table(rows=num_rows, cols=num_cols)
        table2 = graphic_frame.table

        self.fill_table_colour(table2,0,num_rows, num_cols,self.RGB_white)
        #first row
        self.fill_table_colour(table2,0,1, num_cols,self.RGB_light_blue)  

        for row in range(num_rows):              
            cell_start  = table2.cell(row, 0)
            cell_end    = table2.cell(row, 2)
            cell_start.merge(cell_end)

        #fill table with data
        cell2=table1.cell(0, 0)
        cell2.text="BASE PEOPLE & CONVERSATIONS"
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        index=0
        for i,row in data_df1.iterrows(): 
            #logging.debug(f" ---> row= {row}")
            #logging.debug(f" ---> index= {index}")               

            table1.cell(index+1, 0).text=row["Type"] 
 
            cell=table1.cell(index+1, 3)
            cell.text=str(row["MeasureValue"])
            cell.alignment = PP_ALIGN.CENTER 
            index+=1

        cell1=table1.cell(0, 3)
        cell1.text=row["MeasureType"]
        cell1.alignment = PP_ALIGN.CENTER 
        cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        #fill table2 with data
        index=0
        for i,row in data_df2.iterrows(): 
            #logging.debug(f" ---> row= {row}")
            #logging.debug(f" ---> index= {index}")       
            table2.cell(index+1,0).text=row["Type"] 
 
            cell=table2.cell(index+1, 3)
            cell.text=str(row["MeasureValue"])
            cell.alignment = PP_ALIGN.CENTER 
            index+=1

        cell2=table2.cell(0, 0)
        cell2.text=row["Section"].upper()
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell1=table2.cell(0, 3)
        cell1.text=row["MeasureType"]
        cell1.alignment = PP_ALIGN.CENTER 
        cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)        

    def create_slide_layout3(self,prs,data_df1,data_df2,data_df3,data_df4):
        logging.debug(f" ---> data_df1={data_df1}")   
        logging.debug(f" ---> data_df1={data_df2}")
        logging.debug(f" ---> data_df1={data_df3}")
        logging.debug(f" ---> data_df1={data_df4}") 

        #####slide3 - two content slide - layout 3
        slide_layout_slide3 = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout_slide3)

        #get ph in template
        title = slide.shapes.title

        table1_ph = slide.placeholders[18]
        table2_ph = slide.placeholders[19]

        #fill ph with content
        title.text = "Tables with merged header"
        
        num_rows=len(data_df1)+2
        num_cols=10
        graphic_frame = table1_ph.insert_table(rows=num_rows, cols=num_cols)
        table1= graphic_frame.table 
        self.fill_table_colour(table1,0,num_rows, num_cols,self.RGB_white)
        #first 2 rows
        self.fill_table_colour(table1,0,1, num_cols,self.RGB_biege)  
        self.fill_table_colour(table1,1,2, num_cols,self.RGB_light_blue) 

        cell_start  = table1.cell(0, 0)
        cell_end    = table1.cell(0, num_cols-1)
        cell_start.merge(cell_end)

        for row in range(1,num_rows):                 
            cell_start  = table1.cell(row, 0)
            cell_end    = table1.cell(row, 3)
            cell_start.merge(cell_end)    

            cell_start  = table1.cell(row, 4)
            cell_end    = table1.cell(row, 5)
            cell_start.merge(cell_end) 

            cell_start  = table1.cell(row, 6)
            cell_end    = table1.cell(row, 8)
            cell_start.merge(cell_end)    

        num_rows=len(data_df3)+2
        num_cols=10
        graphic_frame = table2_ph.insert_table(rows=num_rows, cols=num_cols)
        table2 = graphic_frame.table 
        self.fill_table_colour(table2,0,num_rows, num_cols,self.RGB_white)
        #first row
        self.fill_table_colour(table2,0,1, num_cols,self.RGB_biege) 
        self.fill_table_colour(table2,1,2, num_cols,self.RGB_light_blue) 
              
        cell_start  = table2.cell(0, 0)
        cell_end    = table2.cell(0, num_cols-1)
        cell_start.merge(cell_end) 

        for row in range(1,num_rows):                
            cell_start  = table2.cell(row, 0)
            cell_end    = table2.cell(row, 3)
            cell_start.merge(cell_end) 

            cell_start  = table2.cell(row, 4)
            cell_end    = table2.cell(row, 5)
            cell_start.merge(cell_end) 

            cell_start  = table2.cell(row, 6)
            cell_end    = table2.cell(row, 8)
            cell_start.merge(cell_end)             

        #fill table1 with data
        #data_df1
        index=0
        for i,row in data_df1.iterrows(): 
            #logging.debug(f" ---> row= {row}")
            #logging.debug(f" ---> index= {index}")               

            table1.cell(index+2,0).text=row["Type"] 
 
            cell=table1.cell(index+2, 4)
            cell.text=str(row["MeasureValue"])
            cell.alignment = PP_ALIGN.CENTER 
            index+=1

        cell2=table1.cell(0, 0)
        cell2.text=row["Section"]
        cell2.alignment = PP_ALIGN.CENTER
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell2=table1.cell(1, 0)
        cell2.text=row["Segment"]
        cell2.alignment = PP_ALIGN.CENTER
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell1=table1.cell(1, 4)
        cell1.text=row["MeasureType"]
        cell1.alignment = PP_ALIGN.CENTER 
        cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)   

        #data_df2
        index=0
        for i,row in data_df2.iterrows(): 
            #logging.debug(f" ---> row= {row}")
            #logging.debug(f" ---> index= {index}")       
            table1.cell(index+2,6).text=row["Type"] 
 
            cell=table1.cell(index+2, 9)
            cell.text=str(row["MeasureValue"])
            cell.alignment = PP_ALIGN.CENTER 
            index+=1

        cell2=table1.cell(0, 0)
        cell2.text=row["Section"]
        cell2.alignment = PP_ALIGN.CENTER
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell2=table1.cell(1, 6)
        cell2.text=row["Segment"]
        cell2.alignment = PP_ALIGN.CENTER
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell1=table1.cell(1, 9)
        cell1.text=row["MeasureType"]
        cell1.alignment = PP_ALIGN.CENTER 
        cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)   

        #fill table2 with data       
        #data_df3
        index=0
        for i,row in data_df3.iterrows(): 
            #logging.debug(f" ---> row= {row}")
            #logging.debug(f" ---> index= {index}")                

            table2.cell(index+2,0).text=row["Type"] 
 
            cell=table2.cell(index+2, 4)
            cell.text=str(row["MeasureValue"])
            cell.alignment = PP_ALIGN.CENTER 
            index+=1

        cell2=table2.cell(0, 0)
        cell2.text=row["Section"]
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell2=table2.cell(1, 0)
        cell2.text=row["Segment"]
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell1=table2.cell(1, 4)
        cell1.text=row["MeasureType"]
        cell1.alignment = PP_ALIGN.CENTER 
        cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)   

        #data_df4
        index=0
        for i,row in data_df4.iterrows(): 
            #logging.debug(f" ---> row= {row}")
            #logging.debug(f" ---> index= {index}")       

            table2.cell(index+2,6).text=row["Type"] 
 
            cell=table2.cell(index+2, 9)
            cell.text=str(row["MeasureValue"])
            cell.alignment = PP_ALIGN.CENTER 
            index+=1

        cell2=table2.cell(0, 0)
        cell2.text=row["Section"]
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell2=table2.cell(1, 6)
        cell2.text=row["Segment"]
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell1=table2.cell(1, 9)
        cell1.text=row["MeasureType"]
        cell1.alignment = PP_ALIGN.CENTER 
        cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)   

    def create_slide4(self,prs,data_df1):    
        logging.debug(f" ---> data_df1={data_df1}")                                       
        ######slide4: - two content slide - layout 3
        slide_layout_slide3 = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout_slide3)

        #get ph in template
        title = slide.shapes.title
        table1_ph = slide.placeholders[18]

        #fill ph with content
        title.text = "The last slide!"

        num_rows=len(data_df1)+2
        num_cols=10
        graphic_frame = table1_ph.insert_table(rows=num_rows, cols=num_cols)
        table1 = graphic_frame.table 

        self.fill_table_colour(table1,0,num_rows, num_cols,self.RGB_white)
        #first row
        self.fill_table_colour(table1,0,1, num_cols,self.RGB_biege)  
        self.fill_table_colour(table1,1,2, num_cols,self.RGB_light_blue)         

        cell_start  = table1.cell(0, 0)
        cell_end    = table1.cell(0, 9)
        cell_start.merge(cell_end)         

        for row in range(1,num_rows):              
            cell_start  = table1.cell(row, 0)
            cell_end    = table1.cell(row, 6)
            cell_start.merge(cell_end)  

            cell_start  = table1.cell(row, 7)
            cell_end    = table1.cell(row, 9)
            cell_start.merge(cell_end) 

        #fill table1 with data
        index=0
        for i,row1 in data_df1.iterrows(): 
            #logging.debug(f" ---> row1= {row1}")
            #logging.debug(f" ---> index= {index}")                        

            table1.cell(index+2,0).text=row1["Type"] 
 
            cell=table1.cell(index+2, 7)
            cell.text=str(row1["MeasureValue"])
            cell.alignment = PP_ALIGN.CENTER 
            index+=1

        cell2=table1.cell(0, 0)
        cell2.text=row1["Section"]
        cell2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        cell1=table1.cell(1, 7)
        cell1.text=row1["MeasureType"]
        cell1.alignment = PP_ALIGN.CENTER 
        cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)   

    def create_slide_chart(self,prs,data_df1,chart):    
        logging.debug(f" ---> data_df1={data_df1}")

        #data_df1.drop(columns=['PeriodKey','Category','Subcategory','Country','Segment','Section','Content','MeasureType'])
        #logging.debug(f" ---> data_df1={data_df1}")

        self.create_chart(data_df1,chart)

        logging.debug(f" ---> data_df1={data_df1}")

        ######slide4: - content slide - layout 2
        slide_layout_slide = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout_slide)

        #get ph in template
        title = slide.shapes.title
        #content = slide.placeholders[15]

        picture_placeholder = slide.placeholders[15]

        #pic = picture_placeholder.insert_picture(chart)
        pic = slide.shapes.add_picture(chart, Inches(2), Inches(1.5))

        #fill ph with content
        title.text = "The last slide!"             

    def fill_table_colour(self,table_x,start_row,end_row, num_cols,rgb_color):
        #logging.debug(f" ---> num_rows={end_row}")

        for row in range(start_row,end_row):  
            for col in range(num_cols): 
                cell  = table_x.cell(row, col)              
                cell.fill.solid()
                # set foreground (fill) color to a specific RGB color
                cell.fill.fore_color.rgb = RGBColor(rgb_color[0], rgb_color[1], rgb_color[2]) 
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 255)      

    #===========================================================================================
    # Load Data from csv files from local
    #===========================================================================================
    def load_data(self,dir,file_type="csv",move_file=False,skip_rows=0,skip_footer=0):
        #logging.debug(f'load_data---> dir={dir}')

        err_msg=None
        #get from local csv files
        df_dict={}
  
        with os.scandir(dir) as fileList:
            for file in fileList:
                full_path=f'{dir}\{file.name}'
                #logging.debug(f'---> full_path={full_path}')
                #logging.debug(f'---> path.isfile(full_path)={path.isfile(full_path)}') 
                #logging.debug(f'---> os.path.splitext(full_path)[1].lower()={os.path.splitext(full_path)[1].lower()==".csv"}')                               
                #only load csv file
                if (path.isfile(full_path) and os.path.splitext(full_path)[1].lower()==f'.{file_type}'):       
                    df = pd.read_csv(full_path, engine='python', skiprows=skip_rows, skipfooter=skip_footer)
                    #pd.read_csv(full_path, skiprows=17, skipfooter=1)

                    #clean data
                    df.fillna('', inplace=True) #e.g website can be blank

                    #store to list
                    df_dict[file.name]=df

        if df_dict:
            if move_file:
                #move processed files to processed directory
                self.move_files(self.CSV_DIR,self.PROCESSED_CSV_DIR)
                err_msg=None
        else:
            err_msg=f"No {file_type} files to process"

        return (df_dict,err_msg)

    #===========================================================================================
    # #V2 Move all processed files to processed folder in Local dir
    #============================================================================================    
    def move_files(self,from_dir,to_dir):
        #logging.debug(f"move_files ---> from_dir= {from_dir}")   
        
        with os.scandir(from_dir) as fileList:
            for file in fileList:
                full_path=f'{from_dir}\{file.name}'
                #only move csv file
                if (path.isfile(full_path) and os.path.splitext(full_path)[1].lower()==".csv"): 
                    logging.debug(f'Moving {file.name}')

                    original_file=fr'{from_dir}\{file.name}'
                    new_file=fr'{to_dir}\{file.name}'

                    os.rename(original_file,new_file)

#main
log_format="[ %(asctime)s - %(levelname)s - %(threadName)s - (%(name)s - %(filename)s - %(funcName)s(), line %(lineno)d)]: %(message)s"
logging.basicConfig(filename='q2c.log', filemode='w', level=logging.DEBUG,format=log_format)

logging.info("MAIN ---> START")

print(os.getcwd())
cwd=os.getcwd()

Q_DIR=rf"{cwd}\Q2"
CSV_DIR=rf"{Q_DIR}\2c"
PROCESSED_CSV_DIR=rf"{CSV_DIR}\processed"
OUTPUT_CSV_DIR=rf"{CSV_DIR}\output"

q2c=Q2c()

(df_list,err_msg)=q2c.load_data(CSV_DIR)

input_file=rf"{Q_DIR}\template.pptx"
output_file="output_chart.pptx"

for key, df in df_list.items():
    #logging.debug(f"---> key= {key}")
    #logging.debug(f"---> df= {df}")
    q2c.create_ppt(input_file, output_file,df, "pie.png")

logging.info("MAIN ---> END")